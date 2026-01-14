"""Generate PowerPoint presentations from slide content."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os

class PPTXGenerator:
    """Generate PowerPoint presentations with professional styling."""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Professional black and white color scheme
        self.primary_color = RGBColor(0, 0, 0)          # Black
        self.secondary_color = RGBColor(64, 64, 64)     # Dark gray
        self.accent_color = RGBColor(128, 128, 128)     # Medium gray
        self.text_color = RGBColor(32, 32, 32)          # Near black
        self.bg_color = RGBColor(255, 255, 255)         # White
        
    def add_title_slide(self, title: str, subtitle: str = "", authors: str = ""):
        """Add a professionally styled title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout for custom design
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # Add colored header bar
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1.5)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.8),
                Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = self.secondary_color
            p.alignment = PP_ALIGN.CENTER
        
        # Add authors/date
        if authors:
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(6),
                Inches(8), Inches(0.8)
            )
            author_frame = author_box.text_frame
            p = author_frame.paragraphs[0]
            p.text = authors
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            p.alignment = PP_ALIGN.CENTER
        
    def add_content_slide(self, title: str, bullets: list, notes: str = ""):
        """Add a professionally styled content slide with bullets."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # Add header bar
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(0.8)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add accent line
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1),
            Inches(0.1), Inches(5.5)
        )
        accent_fill = accent_line.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = self.accent_color
        accent_line.line.fill.background()
        
        # Add content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.2),
            Inches(8.5), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        # Add bullets with better formatting
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(12)
            p.space_after = Pt(6)
            p.line_spacing = 1.2
            
        # Add footer with slide number
        footer_box = slide.shapes.add_textbox(
            Inches(9), Inches(7.2),
            Inches(0.8), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        p = footer_frame.paragraphs[0]
        p.text = str(len(self.prs.slides))
        p.font.size = Pt(12)
        p.font.color.rgb = self.text_color
        p.alignment = PP_ALIGN.RIGHT
        
        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
            
    def add_section_slide(self, section_title: str):
        """Add a section divider slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.primary_color
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide_with_image(self, title: str, bullets: list, image_path: str = None):
        """Add a content slide with bullets and an optional image."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # Add header bar
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(0.8)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        # Add title with word wrap
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.1),
            Inches(9), Inches(0.65)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)  # Smaller to fit long titles
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Determine layout based on image availability
        if image_path and os.path.exists(image_path):
            # Two-column layout: text on left, image on right
            content_width = 4.5
            image_left = Inches(5.2)
            image_width = Inches(4.3)
        else:
            # Full width for text
            content_width = 8.5
            image_path = None
        
        # Add accent line
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1),
            Inches(0.1), Inches(5.5)
        )
        accent_fill = accent_line.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = self.accent_color
        accent_line.line.fill.background()
        
        # Add content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.2),
            Inches(content_width), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        # Add bullets
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(10)
            p.space_after = Pt(6)
            p.line_spacing = 1.2
        
        # Add image if available
        if image_path:
            try:
                pic = slide.shapes.add_picture(
                    image_path,
                    image_left,
                    Inches(1.5),
                    width=image_width
                )
            except Exception as e:
                print(f"Could not add image: {e}")
        
        # Add footer
        footer_box = slide.shapes.add_textbox(
            Inches(9), Inches(7.2),
            Inches(0.8), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        p = footer_frame.paragraphs[0]
        p.text = str(len(self.prs.slides))
        p.font.size = Pt(12)
        p.font.color.rgb = self.text_color
        p.alignment = PP_ALIGN.RIGHT
    
    def add_qa_slide(self, title: str, qa_pairs: list):
        """Add a Q&A slide with questions and answers."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # Add header bar
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(0.8)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.accent_color
        header_shape.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add Q&A content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        for i, qa in enumerate(qa_pairs):
            if i > 0:
                p = text_frame.add_paragraph()
                p.space_before = Pt(18)
            else:
                p = text_frame.paragraphs[0]
            
            # Split Q and A
            if '\nA:' in qa:
                q_part, a_part = qa.split('\nA:', 1)
                q_text = q_part.replace('Q:', '').strip()
                a_text = a_part.strip()
                
                # Add question
                p.text = f"❓ {q_text}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self.primary_color
                
                # Add answer
                p_answer = text_frame.add_paragraph()
                p_answer.text = f"   {a_text}"
                p_answer.font.size = Pt(14)
                p_answer.font.color.rgb = self.text_color
                p_answer.space_before = Pt(6)
                p_answer.space_after = Pt(12)
        
    def save(self, filename: str):
        """Save the presentation."""
        self.prs.save(filename)
        

def parse_slide_content(slide_text: str) -> dict:
    """Parse slide content from text format."""
    lines = slide_text.strip().split('\n')
    title = lines[0].replace('**', '').replace('*', '').strip()
    
    bullets = []
    visual = ""
    
    for line in lines[1:]:
        line = line.strip()
        if line.startswith('*') or line.startswith('-') or line.startswith('•'):
            # Remove bullet markers and clean
            bullet = re.sub(r'^[\*\-•]\s*', '', line).strip()
            if bullet and not bullet.lower().startswith('visual:'):
                bullets.append(bullet)
        elif 'visual:' in line.lower():
            visual = line.split(':', 1)[1].strip() if ':' in line else ""
            
    return {
        'title': title,
        'bullets': bullets,
        'visual': visual
    }


def generate_pptx_from_blueprint(blueprint_text: str, output_path: str, paper_title: str = "Research Paper", 
                                 extracted_images: list = None):
    """Generate PowerPoint from slide blueprint text with extracted images."""
    generator = PPTXGenerator()
    
    # Add title slide
    generator.add_title_slide(paper_title, "AI-Generated Presentation", "Created by Multi-Agent System")
    
    # Extract presenter notes section
    presenter_notes_section = ""
    if "**Presenter Notes:**" in blueprint_text:
        parts = blueprint_text.split("**Presenter Notes:**")
        blueprint_text = parts[0]
        presenter_notes_section = parts[1] if len(parts) > 1 else ""
    
    # Split into slides - handle both formats
    # Try numbered format first: "1. **Title**"
    slides = re.split(r'\n\d+\.\s+\*\*', blueprint_text)
    
    # If that didn't work, try section format: "=== TITLE ==="
    if len(slides) <= 1:
        slides = re.split(r'===\s+([^=]+)\s+===', blueprint_text)
    
    # Process content slides with images
    image_index = 0
    for i, slide_text in enumerate(slides[1:]):  # Skip first empty split
        # For section format, odd indices are titles, even are content
        if '===' in blueprint_text and i % 2 == 0:
            title = slide_text.strip()
            content = slides[i + 2] if i + 2 < len(slides) else ""
            bullets = [line.strip('- ').strip() for line in content.split('\n') if line.strip().startswith('-')]
        else:
            slide_data = parse_slide_content(slide_text)
            title = slide_data['title']
            bullets = slide_data['bullets']
        
        if bullets:
            # Add image if available
            image_path = None
            if extracted_images and image_index < len(extracted_images):
                image_path = extracted_images[image_index]['path']
                image_index += 1
            
            generator.add_content_slide_with_image(
                title,
                bullets,
                image_path
            )
    
    # Add Q&A slide at the end
    if presenter_notes_section:
        qa_content = extract_qa_from_notes(presenter_notes_section)
        if qa_content:
            generator.add_qa_slide("Questions & Discussion", qa_content)
    
    generator.save(output_path)
    return output_path


def extract_qa_from_notes(notes_text: str) -> list:
    """Extract questions and answers from presenter notes."""
    qa_pairs = []
    
    # Find all question-answer pairs
    lines = notes_text.split('\n')
    current_question = None
    
    for line in lines:
        line = line.strip()
        if '?' in line and ('What' in line or 'How' in line or 'Why' in line):
            # This is a question
            question = line.replace('+ ', '').replace('*', '').strip()
            if question and len(question) < 200:  # Reasonable length
                current_question = question
        elif current_question and line.startswith('+ ') and 'answer' not in line.lower():
            # This might be an answer
            answer = line.replace('+ ', '').strip()
            if answer and len(answer) < 300:
                qa_pairs.append(f"Q: {current_question}\nA: {answer}")
                current_question = None
    
    # Limit to top 5 Q&A pairs
    return qa_pairs[:5]
